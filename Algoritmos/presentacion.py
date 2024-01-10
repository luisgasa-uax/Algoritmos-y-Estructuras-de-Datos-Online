from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a PowerPoint presentation
prs = Presentation()

# Define title and detailed content for each slide
slides_content = [
    {"title": "Diseño de Algoritmos - Unidad Didáctica 6", "content": ""},

    {"title": "Introducción", "content": "Exploramos el diseño, implementación y análisis de algoritmos, enfocándonos en la eficiencia y la aplicación práctica."},

    {"title": "Integración de Estructuras de Datos y Algoritmos Clásicos", 
     "content": "Análisis de cómo estructuras como listas, pilas, colas y grafos se combinan con algoritmos de ordenación y búsqueda para resolver problemas complejos."},

    {"title": "Algoritmos de Ordenación y Búsqueda", 
     "content": "Estudio detallado de algoritmos como BubbleSort, Insertion Sort, Merge Sort y QuickSort. Comparación de su eficiencia y aplicabilidad en distintos contextos."},

    {"title": "Aplicación Práctica de Algoritmos", 
     "content": "Ejemplos de cómo se utilizan estos algoritmos en la vida real, como en sistemas de gestión de bases de datos y optimización de procesos."},

    {"title": "Algoritmos Avanzados", 
     "content": "Introducción a algoritmos más complejos, como Dijkstra para rutas mínimas y Floyd-Warshall para análisis de grafos."},

    {"title": "Programación Dinámica y Memorización", 
     "content": "Discusión sobre técnicas de optimización para resolver problemas complejos, con ejemplos como el problema de la mochila y algoritmos recursivos con memorización."},

    {"title": "Conclusión", 
     "content": "Resumen de conceptos clave, destacando la importancia de seleccionar el algoritmo adecuado para cada situación específica."},

    {"title": "Referencias", 
     "content": "Fuentes y materiales de referencia utilizados para el desarrollo de esta unidad didáctica."}
]


# Function to add a slide
def add_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]  # Using layout 1 for title and content
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    body_placeholder = slide.placeholders[1]
    
    title_placeholder.text = title
    
    tf = body_placeholder.text_frame
    tf.text = content

    for paragraph in tf.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(18)
            run.font.color.rgb = RGBColor(0, 0, 0)

# Add slides to the presentation
for slide in slides_content:
    add_slide(prs, slide["title"], slide["content"])

# Save the presentation
pptx_file = "UD6_Diseno_De_Algoritmos_Presentacion.pptx"
prs.save(pptx_file)

pptx_file

